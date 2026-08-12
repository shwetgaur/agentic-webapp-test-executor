from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

path = Path(r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (3).pptx")
prs = Presentation(path)
print("size", prs.slide_width, prs.slide_height)
print("slides", len(prs.slides))
print("layouts", [(i, l.name) for i, l in enumerate(prs.slide_layouts)])

for si, slide in enumerate(prs.slides):
    print(f"\n{'='*60}\nSLIDE {si+1}")
    for shape in slide.shapes:
        print(f"  [{shape.shape_type}] {shape.name} L={shape.left} T={shape.top} W={shape.width} H={shape.height}")
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                t = p.text
                if t.strip():
                    print(f"    P{pi}: {t!r}")
        if shape.has_table:
            table = shape.table
            print(f"    TABLE {len(table.rows)}x{len(table.columns)}")
            for r in table.rows:
                cells = [c.text.replace("\n", " | ") for c in r.cells]
                print("     |", " || ".join(cells))
