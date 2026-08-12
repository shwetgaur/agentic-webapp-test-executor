"""Generate B.Tech CA1 presentation matching college template format."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

ROOT = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor")
ASSETS = ROOT / "docs" / "ca1_assets"
OUT = ROOT / "docs" / "DS1_Agentic_WebApp_Test_Executor_CA1_Presentation.pptx"

# Standard widescreen used by many SIT decks; template PDF looks 4:3-ish but
# college PPT templates often use 13.33x7.5. Keep close to inspected PBL size.
# PBL template was 9144000 x 6858000 EMU ~= 10" x 7.5" (4:3)
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(6858000)

RED = RGBColor(0xC4, 0x1E, 0x3A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x1E, 0x7A, 0x46)
ORANGE = RGBColor(0xC6, 0x5A, 0x12)
PURPLE = RGBColor(0x5B, 0x2C, 0x6F)

DATE = "28-07-2026"
DEPT = "Department of Artificial Intelligence & Machine Learning"

TITLE_LOGO = ASSETS / "sit_logo_title.png"
HEADER_LOGO = ASSETS / "sit_logo_header.png"


def set_run_font(run, *, size=18, bold=False, color=DARK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.get_or_add_latin()
    latin.set("typeface", name)


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color, name=font)
    return box


def add_bullets(slide, left, top, width, height, items, *, size=16, color=DARK, bold_first=False, level_gap=0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item
        set_run_font(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_footer(slide, number: int):
    # left date
    add_textbox(slide, Inches(0.35), Inches(7.05), Inches(2.2), Inches(0.3), DATE, size=11, color=GRAY)
    # center dept
    add_textbox(
        slide,
        Inches(2.4),
        Inches(7.05),
        Inches(5.5),
        Inches(0.3),
        DEPT,
        size=11,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    # right number
    add_textbox(slide, Inches(8.7), Inches(7.05), Inches(0.9), Inches(0.3), str(number), size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def add_header_logo(slide):
    if HEADER_LOGO.exists():
        slide.shapes.add_picture(str(HEADER_LOGO), Inches(7.15), Inches(0.12), height=Inches(0.55))


def add_title_bar(slide, title: str):
    add_header_logo(slide)
    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(6.6), Inches(0.5), title, size=26, bold=True, color=DARK, font="Georgia")


def rounded_card(slide, left, top, width, height, fill, text, *, text_size=12, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run_font(run, size=text_size, bold=True, color=text_color)
    shape.text_frame.auto_size = None
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    return shape


def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.35), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()
    return shape


def new_content_slide(prs, title: str, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title_bar(slide, title)
    add_footer(slide, number)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ========== 1 TITLE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4), "B. Tech. Project AY 2026-27", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), "CA 1 Presentation", size=22, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(1.15), Inches(9), Inches(0.3), "Project ID: DS 1", size=16, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(0.4),
        Inches(1.5),
        Inches(9.2),
        Inches(0.55),
        "Title of the Project: Agentic Web-App Test Executor\n(Domain: Quality Engineering)",
        size=16,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if TITLE_LOGO.exists():
        s.shapes.add_picture(str(TITLE_LOGO), Inches(2.7), Inches(2.2), width=Inches(4.6))
    add_textbox(s, Inches(0.5), Inches(4.35), Inches(9), Inches(0.3), "Name of the Guide: ________________________", size=15, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.7), Inches(9), Inches(0.3), "Name of the Industry Mentor: Dassault Systemes (ENOVIA) / ________________", size=14, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.2), Inches(8.5), Inches(1.5),
                "Group Members:\n1. ________________ (PRN: ____________)\n2. ________________ (PRN: ____________)\n3. ________________ (PRN: ____________)\n4. ________________ (PRN: ____________)",
                size=14)
    add_footer(s, 1)

    # ========== 2 OUTLINE ==========
    s = new_content_slide(prs, "Outline of the Presentation", 2)
    outline = [
        "1) Names of all the group members and their work distribution",
        "2) Introduction",
        "3) Problem Statement",
        "4) Objectives of the Project",
        "5) Project plan with timeline",
        "6) Literature Review",
        "7) Gap in the Research/ Technology/ Methodology",
        "8) Description of the proposed solution",
        "9) Requirement Analysis",
        "10) Technology Stack",
        "11) Design",
        "12) Development/ Implementation",
        "13) Testing & Debugging",
        "14) Project Outcome",
        "15) Conclusion",
        "16) References",
        "17) A photograph of the presentation after due permission from the guide.",
    ]
    # two columns
    mid = 9
    add_bullets(s, Inches(0.5), Inches(0.95), Inches(4.5), Inches(5.8), outline[:mid], size=14)
    add_bullets(s, Inches(5.1), Inches(0.95), Inches(4.5), Inches(5.8), outline[mid:], size=14)

    # ========== 3 WORK DISTRIBUTION ==========
    s = new_content_slide(prs, "1) Group Members & Work Distribution", 3)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Member A — Executor: Playwright actions, assertions, screenshots, demo flows",
            "Member B — AI Agent: text→JSON step parser, prompts/LLM, architecture, paper draft",
            "Member C — Platform: FastAPI, DB, dashboard, report storage/export, notify plumbing",
            "Member D — Notify + Docs: team ownership map, failure alerts, lit review, report, logbook, slides",
            "Shared (ALL): weekly guide meetings, objectives freeze, GitHub hygiene, final demo",
            "Note: End-to-end pipeline already scaffolded so any member can integrate without blockers",
        ],
        size=15,
    )

    # ========== 4 INTRODUCTION ==========
    s = new_content_slide(prs, "2) Introduction", 4)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Industry partner: Dassault Systemes (ENOVIA) — Quality Engineering domain",
            "Web application testing is critical for release quality, but manual QA is slow and costly",
            "Traditional automation scripts (Selenium/Playwright code) are brittle when UI changes",
            "Agentic AI can interpret plain-text test steps and execute them with verification",
            "This project builds an Agentic Web-App Test Executor with pass/fail documentation and scrum-style failure notification to owning teams",
        ],
        size=16,
    )

    # ========== 5 PROBLEM STATEMENT ==========
    s = new_content_slide(prs, "3) Problem Statement", 5)
    add_textbox(
        s,
        Inches(0.5),
        Inches(1.0),
        Inches(9.0),
        Inches(1.3),
        "Uses basic text steps as input to automatically replay user actions on a web app. It can follow those steps, execute the flow, and verify whether the app behaves as expected.",
        size=16,
        bold=True,
    )
    # challenge cards
    rounded_card(s, Inches(0.5), Inches(2.6), Inches(2.8), Inches(1.5), RED, "Challenge\nManual testing is slow,\ninconsistent & expensive", text_size=13)
    rounded_card(s, Inches(3.6), Inches(2.6), Inches(2.8), Inches(1.5), ORANGE, "Challenge\nScripts break often\nwhen UI changes", text_size=13)
    rounded_card(s, Inches(6.7), Inches(2.6), Inches(2.8), Inches(1.5), BLUE, "Need\nAI executor + reports\n+ team notify on fail", text_size=13)
    add_bullets(
        s,
        Inches(0.5),
        Inches(4.4),
        Inches(9),
        Inches(2.2),
        [
            "Mentor additions: (1) document what passed/failed with evidence; (2) notify responsible team like a scrum master",
            "Suggested tools (3DS): Playwright, Testim etc. | Sample apps: ClickUp / Zoho (Phase-2 if access allowed)",
        ],
        size=14,
    )

    # ========== 6 OBJECTIVES ==========
    s = new_content_slide(prs, "4) Objectives of the Project", 6)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Interpret basic text test steps into structured executable actions",
            "Execute actions on a web app using Playwright automation",
            "Verify behavior via assertions (text, URL, visibility)",
            "Generate structured pass/fail documentation with screenshots/logs",
            "Notify owning team on failure (scrum-style alert + ticket id)",
            "Reduce script brittleness via AI-assisted parsing / locator recovery",
            "Evaluate on demo apps and produce college + Dassault documentation/paper draft",
        ],
        size=15,
    )

    # ========== 7 TIMELINE ==========
    s = new_content_slide(prs, "5) Project Plan with Timeline", 7)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.9),
        Inches(9.1),
        Inches(5.8),
        [
            "Week 1–2: Kickoff, mentor prep, role lock, repo setup",
            "Week 3: Literature (5+), architecture, step/report schemas freeze",
            "Week 4–6: Executor MVP + parser integration (3 end-to-end flows)",
            "Week 7–8: Pass/fail reports + dashboard",
            "Week 9: Notify agent (team routing + ticket) — mentor feature",
            "Week 10–11: Batch suites, basic self-heal, evaluation metrics",
            "Week 12: Demo freeze, README polish",
            "Week 13–15: Report chapters, slides, poster, paper draft, final packaging",
        ],
        size=15,
    )

    # ========== 8 LITERATURE ==========
    s = new_content_slide(prs, "6) Literature Review", 8)
    add_bullets(
        s,
        Inches(0.4),
        Inches(0.9),
        Inches(9.2),
        Inches(5.8),
        [
            "1. Playwright Test Agents (Planner/Generator/Healer) — Microsoft docs: agentic plan→test→heal pipeline",
            "2. Tricentis Testim — commercial AI testing with smart locators / NL authoring (named by 3DS)",
            "3. Nass et al. (2024), STVR — LLM improves web element localization when locators fail",
            "4. ICECIT 2025 — Generative AI self-healing for Selenium; LLMs vs rule-based repair",
            "5. JAIGS 2025 — RL (PPO) + dynamic XPath self-healing in Playwright",
            "6. Bonus baseline: Robula+ — robust locator generation (pre-AI foundation)",
        ],
        size=14,
    )

    # ========== 9 GAP ==========
    s = new_content_slide(prs, "7) Gap in Research / Technology / Methodology", 9)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Existing tools generate code, heal locators, or offer closed commercial NL testing",
            "Few academic prototypes unify: plain-text execution + evidence-backed pass/fail docs + team failure routing",
            "Reporting is often an afterthought; ownership-aware notify (scrum-master style) is rarely first-class",
            "Our contribution: open, explainable Quality Engineering workflow combining all three",
        ],
        size=16,
    )

    # ========== 10 PROPOSED SOLUTION ==========
    s = new_content_slide(prs, "8) Description of the Proposed Solution", 10)
    # pipeline diagram
    boxes = [
        (0.35, "Text Steps", RED),
        (2.15, "Step Parser\n(Rules/LLM)", BLUE),
        (3.95, "Playwright\nExecutor", GREEN),
        (5.75, "Assertions\nPass/Fail", ORANGE),
        (7.55, "Report +\nNotify", PURPLE),
    ]
    y = Inches(1.15)
    for i, (x, label, color) in enumerate(boxes):
        rounded_card(s, Inches(x), y, Inches(1.65), Inches(1.05), color, label, text_size=11)
        if i < len(boxes) - 1:
            arrow_right(s, Inches(x + 1.68), Inches(1.5))
    add_bullets(
        s,
        Inches(0.45),
        Inches(2.5),
        Inches(9.1),
        Inches(4.2),
        [
            "Input: numbered plain-text steps (login, fill, click, verify...)",
            "Parser maps steps to locked JSON action schema (goto/fill/click/assert/...)",
            "Executor runs actions in Chromium via Playwright; captures failure screenshots",
            "Report writer stores JSON + Markdown pass/fail documentation",
            "Notify agent maps module→owning team and raises ticket-style alert on FAIL",
            "Phase-1 demo target: Sauce Demo (public). ClickUp/Zoho later if mentors grant access",
        ],
        size=14,
    )

    # ========== 11 REQUIREMENTS ==========
    s = new_content_slide(prs, "9) Requirement Analysis", 11)
    add_textbox(s, Inches(0.5), Inches(0.95), Inches(4.4), Inches(0.35), "Functional Requirements", size=16, bold=True, color=RED)
    add_bullets(
        s,
        Inches(0.5),
        Inches(1.35),
        Inches(4.5),
        Inches(5.2),
        [
            "Accept plain-text / JSON suites",
            "Execute browser actions automatically",
            "Assert expected UI outcomes",
            "Produce step-level pass/fail reports",
            "Route failures to owning team",
            "CLI + API entry points",
        ],
        size=14,
    )
    add_textbox(s, Inches(5.2), Inches(0.95), Inches(4.4), Inches(0.35), "Non-Functional Requirements", size=16, bold=True, color=RED)
    add_bullets(
        s,
        Inches(5.2),
        Inches(1.35),
        Inches(4.5),
        Inches(5.2),
        [
            "Reliability of step execution",
            "Evidence retention (screenshots)",
            "Extensible schemas (v1 locked)",
            "Maintainable modular codebase",
            "Demoable within minutes",
            "Secure handling of credentials (.env)",
        ],
        size=14,
    )

    # ========== 12 TECH STACK ==========
    s = new_content_slide(prs, "10) Technology Stack", 12)
    stack = [
        (0.5, "Automation", "Playwright\nChromium", BLUE),
        (2.5, "Backend", "Python\nFastAPI", GREEN),
        (4.5, "AI Layer", "LLM / Rules\nParser", PURPLE),
        (6.5, "Data", "JSON/YAML\nSQLite later", ORANGE),
        (8.3, "Notify", "Console /\nSlack/Email", RED),
    ]
    for x, title, body, color in stack:
        add_textbox(s, Inches(x), Inches(1.1), Inches(1.7), Inches(0.35), title, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        rounded_card(s, Inches(x), Inches(1.55), Inches(1.7), Inches(1.3), color, body, text_size=12)
    add_bullets(
        s,
        Inches(0.5),
        Inches(3.2),
        Inches(9),
        Inches(3.5),
        [
            "Language: Python 3.12+",
            "Validation: Pydantic models aligned to step/report schemas",
            "Config: team_ownership.yaml for scrum-style routing",
            "Frontend (next): Streamlit/React dashboard for run history",
            "Optional (3DS wishlist): evaluate Testim concepts; implement openly with Playwright",
        ],
        size=15,
    )

    # ========== 13 DESIGN ==========
    s = new_content_slide(prs, "11) Design", 13)
    # layered architecture boxes
    layers = [
        (1.1, RGBColor(0xEE, 0xEE, 0xEE), DARK, "Presentation / API Layer — CLI + FastAPI endpoints"),
        (2.0, RGBColor(0xD6, 0xEA, 0xF8), DARK, "Agent Layer — Plain-text parser (rules now, LLM next)"),
        (2.9, RGBColor(0xD5, 0xF5, 0xE3), DARK, "Execution Layer — Playwright actions + assertions + screenshots"),
        (3.8, RGBColor(0xFC, 0xF3, 0xCF), DARK, "Knowledge / Report Layer — JSON+Markdown pass/fail documentation"),
        (4.7, RGBColor(0xF5, 0xB7, 0xB1), DARK, "Notify Layer — Module→Team map, ticket id, alert dispatch"),
    ]
    for top, fill, tcol, text in layers:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(top), Inches(8.6), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = GRAY
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_run_font(run, size=14, bold=True, color=tcol)
    add_textbox(s, Inches(0.5), Inches(5.6), Inches(9), Inches(0.9), "Fixed contracts: schemas/step_schema.v1.json • schemas/report_schema.v1.json • config/team_ownership.yaml", size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== 14 IMPLEMENTATION ==========
    s = new_content_slide(prs, "12) Development / Implementation", 14)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Repository scaffolded: src/agent, src/executor, src/reporting, src/notify, src/backend",
            "Rule-based parser converts numbered plain-text cases → TestSuite JSON",
            "PlaywrightExecutor runs goto/fill/click/select/assert_* with failure screenshots",
            "Report writer emits JSON + Markdown under data/reports/",
            "NotifyAgent routes FAIL using module ownership map (console/Slack-ready)",
            "Verified smoke runs: login suite PASS (7/7); intentional fail triggers ticket notify",
            "Next: LLM parser, dashboard UI, more sample flows, self-heal heuristics",
        ],
        size=14,
    )

    # ========== 15 TESTING ==========
    s = new_content_slide(prs, "13) Testing & Debugging", 15)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Unit tests: parser step mapping (pytest)",
            "Sample suites: 10 plain-text cases + JSON login suite (Sauce Demo)",
            "Positive path: successful login verification",
            "Negative path: intentional missing-text assert to validate fail report + notify",
            "Artifacts: screenshots on failure, structured step errors, skip-after-fail policy",
            "Debugging aids: headed mode CLI flag, Markdown report for human review",
        ],
        size=15,
    )

    # ========== 16 OUTCOME ==========
    s = new_content_slide(prs, "14) Project Outcome", 16)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "Working prototype of Agentic Web-App Test Executor (industry PS from Dassault)",
            "Pass/fail documentation module with evidence (mentor requirement)",
            "Scrum-style failure notification to owning teams (mentor requirement)",
            "College deliverables: GitHub, report chapters, poster, demo, logbook",
            "Research paper draft jointly with 3DS mentors (as per Dassault proposal)",
            "Aligned to B.Tech SoP outcomes: industry problem + deployable demo path",
        ],
        size=15,
    )

    # ========== 17 CONCLUSION ==========
    s = new_content_slide(prs, "15) Conclusion", 17)
    add_bullets(
        s,
        Inches(0.45),
        Inches(0.95),
        Inches(9.1),
        Inches(5.7),
        [
            "DS 1 addresses a real Quality Engineering need: NL/text-driven web test execution",
            "Combining execution, documentation, and team notify creates a complete QA workflow",
            "Phase-1 foundation is complete and demoable; Phase-2 focuses on LLM intelligence and UI",
            "Clear path for evaluation, publication, and industry mentor collaboration",
        ],
        size=16,
    )

    # ========== 18 REFERENCES ==========
    s = new_content_slide(prs, "16) References", 18)
    add_bullets(
        s,
        Inches(0.4),
        Inches(0.9),
        Inches(9.2),
        Inches(5.8),
        [
            "Dassault Systemes — AI/ML Projects with SIT ENOVIA 2026 proposal (DS 1 brief)",
            "Microsoft Playwright Test Agents documentation — playwright.dev/docs/test-agents",
            "Tricentis Testim — testim.io",
            "Nass, Alégroth, Feldt (2024). Improving Web Element Localization by Using a Large Language Model. STVR.",
            "ICECIT 2025. Generative AI for Self-Healing Selenium Tests. DOI: 10.1109/icecit67774.2025.11451150",
            "JAIGS 2025. Self-Healing Automation with RL (PPO) in Playwright. DOI: 10.60087/jaigs.v5i1.341",
            "Leotta et al. — Robula+ robust web locator generation",
        ],
        size=13,
    )

    # ========== 19 PHOTO ==========
    s = new_content_slide(prs, "17) Photograph of the Presentation", 19)
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(7.0), Inches(4.2))
    frame.fill.solid()
    frame.fill.fore_color.rgb = LIGHT_GRAY
    frame.line.color.rgb = GRAY
    add_textbox(
        s,
        Inches(1.7),
        Inches(3.0),
        Inches(6.6),
        Inches(1.2),
        "Insert group presentation photograph here\nafter due permission from the guide.",
        size=16,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    prs.save(OUT)
    print("Saved:", OUT)
    print("Slides:", len(prs.slides))


if __name__ == "__main__":
    build()
