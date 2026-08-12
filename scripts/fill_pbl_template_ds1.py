"""Clone PBL-2 CA1 template and replace content for DS1 B.Tech project."""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = Path(r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (3).pptx")
OUT = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\DS1_CA1_Presentation_from_PBL_Template.pptx")
OUT_DL = Path(r"C:\Users\shwet\Downloads\DS1_CA1_Presentation_from_PBL_Template.pptx")

DATE = "28-07-2026"


def set_shape_text(shape, lines: list[str]) -> None:
    """Replace paragraph texts; keep extra paragraphs emptied."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Ensure enough paragraphs
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, p in enumerate(tf.paragraphs):
        if i < len(lines):
            # Clear existing runs and write one run with first run's formatting if present
            if p.runs:
                # keep first run formatting, clear others
                p.runs[0].text = lines[i]
                for r in p.runs[1:]:
                    r.text = ""
            else:
                run = p.add_run()
                run.text = lines[i]
        else:
            if p.runs:
                p.runs[0].text = ""
                for r in p.runs[1:]:
                    r.text = ""


def find_shapes_by_text(slide, substr: str):
    hits = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if substr.lower() in full.lower():
                hits.append(shape)
    return hits


def replace_all_dates(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if "17-01-2025" in r.text or r.text.strip() == "Date":
                        r.text = DATE
                    elif "20-07-2026" in r.text:
                        r.text = DATE


def clear_pictures(slide, keep_none=True):
    """Remove picture shapes from a slide (architecture placeholder)."""
    # Collect picture shape elements to remove
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            to_remove.append(shape._element)
    for el in to_remove:
        sp_tree.remove(el)


def set_table_cell(cell, text: str):
    # Replace first paragraph run text; clear rest
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = text
    for p in tf.paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def fill_table(table, rows_data: list[list[str]]):
    for r_idx, row_vals in enumerate(rows_data):
        if r_idx >= len(table.rows):
            break
        for c_idx, val in enumerate(row_vals):
            if c_idx >= len(table.columns):
                break
            set_table_cell(table.cell(r_idx, c_idx), val)


def main():
    shutil.copy2(SRC, OUT)
    prs = Presentation(OUT)

    # ---- SLIDE 1: Title ----
    s = prs.slides[0]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Project Based Learning" in full or "Semester:" in full:
            set_shape_text(shape, ["B. Tech. Project AY 2026-27", "Semester: VII  Batch: 2023-27"])
        elif "PBL-2 Review" in full or "CA-1" in full:
            set_shape_text(shape, ["B.Tech Project Review Presentation (CA-1)"])
        elif "GroupID" in full or "Group ID" in full:
            set_shape_text(shape, ["Problem ID: DS 1  |  Industry: Dassault Systemes (ENOVIA)"])
        elif "Title of the Project" in full:
            set_shape_text(
                shape,
                [
                    "Title of the Project: Agentic Web-App Test Executor "
                    "(Domain: Quality Engineering) — plain-text steps to automated "
                    "web action replay, verification, pass/fail reporting, and team notify"
                ],
            )
        elif "Name of the Guide" in full:
            set_shape_text(shape, ["Name of the Guide: Dr. Geetanjali S.", "Industry Mentor: Dassault Systemes (ENOVIA)"])
        elif "Group Members" in full:
            set_shape_text(
                shape,
                [
                    "Group Members:",
                    "Nidhi Supe (23070126080)",
                    "Prikshit Gaur (23070126094)",
                    "Shwet Gaur (23070126126)",
                    "Manan Khanna (23070126156)",
                ],
            )

    # ---- SLIDE 2: Outline ----
    s = prs.slides[1]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Outline" in full and len(shape.text_frame.paragraphs) <= 2:
            set_shape_text(shape, ["Outline of the Presentation"])
        elif "Names of all the group members" in full or "Introduction" in full:
            set_shape_text(
                shape,
                [
                    "Names of all the group members and their work distribution",
                    "Introduction",
                    "Problem Statement",
                    "Objectives of the Project",
                    "Project plan with timeline",
                    "Literature Review",
                    "Gap in the Research/ Technology/ Methodology",
                    "Description of the proposed solution",
                    "Architectural Diagram (to be inserted)",
                    "Technology Stack",
                    "Work done till now",
                    "Conclusion",
                    "References",
                    "A photograph of the presentation after due permission from the guide.",
                ],
            )

    # ---- SLIDE 3: Work distribution table ----
    s = prs.slides[2]
    for shape in s.shapes:
        if shape.has_text_frame and "work distribution" in "\n".join(p.text for p in shape.text_frame.paragraphs).lower():
            set_shape_text(shape, ["Name of all the group members and their work distribution"])
        if shape.has_table:
            fill_table(
                shape.table,
                [
                    ["Name", "Work Distribution"],
                    [
                        "Nidhi Supe",
                        "QA sample test cases, pass/fail report review, literature support, evaluation notes.",
                    ],
                    [
                        "Prikshit Gaur",
                        "Platform APIs/DB, dashboard & report storage, notify service plumbing, GitHub structure.",
                    ],
                    [
                        "Shwet Gaur",
                        "End-to-end pipeline ownership: step schema, parser, Playwright executor, notify agent, integration.",
                    ],
                    [
                        "Manan Khanna",
                        "Playwright action coverage, assertions, demo scenarios, failure screenshot tooling.",
                    ],
                ],
            )

    # ---- SLIDE 4: Introduction (icon cards + texts) ----
    s = prs.slides[3]
    intro_map = [
        (
            "Cyber threats",
            "Manual web testing is slow and expensive, while release cycles demand faster Quality Engineering feedback.",
        ),
        (
            "Traditional rule-based",
            "Traditional automation scripts are brittle and break frequently when UI structure or selectors change.",
        ),
        (
            "Machine learning",
            "Agentic AI can interpret plain-text test steps and map them to browser actions with verification.",
        ),
        (
            "Most existing systems",
            "Existing tools are fragmented: execution, evidence-backed reporting, and team failure routing are rarely unified.",
        ),
        (
            "SHIELD-NET integrates",
            "Our executor unifies text-step execution, pass/fail documentation with evidence, and scrum-style notify to owning teams.",
        ),
    ]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if full.strip() == "Introduction" or full.strip() == "Introduction ":
            set_shape_text(shape, ["Introduction"])
            continue
        for key, new in intro_map:
            if key in full:
                set_shape_text(shape, [new])
                break

    # ---- SLIDE 5: Problem Statement ----
    s = prs.slides[4]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Problem Statement" in full and len(full) < 40:
            set_shape_text(shape, ["Problem Statement"])
        elif "SHIELD-NET" in full or "explainable AI-driven IDS" in full:
            set_shape_text(
                shape,
                [
                    "DS 1 (Dassault Systemes): Build an Agentic Web-App Test Executor that takes basic text steps as input, "
                    "automatically replays user actions on a web application, executes the intended flow, and verifies whether "
                    "the app behaves as expected — with pass/fail documentation and failure notification to the responsible team."
                ],
            )
    # Keep decorative picture on problem slide (icons from template)

    # ---- SLIDE 6: Objectives ----
    s = prs.slides[5]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Objectives" in full and len(full) < 50:
            set_shape_text(shape, ["Objectives of the project"])
        elif "Detect malicious" in full or "Evaluate IDS" in full or "Provide explainable" in full or "Bridge research" in full:
            set_shape_text(
                shape,
                [
                    "Interpret plain-text test steps into structured executable actions for web flows (login, form fill, navigate, validate).",
                    "Execute and verify flows automatically using Playwright-based browser automation with evidence capture.",
                    "Generate structured pass/fail documentation (step-level status, screenshots, errors) for every run.",
                    "Notify the owning team on failure using a scrum-master-style routing map (module → team → alert/ticket).",
                ],
            )

    # ---- SLIDE 7: Timeline - keep slide chrome; replace old Gantt image with DS1 plan text ----
    s = prs.slides[6]
    for shape in s.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if "Project plan" in full or "timeline" in full.lower():
                set_shape_text(shape, ["Project plan with timeline"])
    clear_pictures(s)
    box = s.shapes.add_textbox(904372, 1555565, 7337181, 4180919)
    set_shape_text(
        box,
        [
            "Week 1–2: Kickoff, mentor prep, repo setup, role lock",
            "Week 3: Literature (≥5), architecture freeze, step/report schemas",
            "Week 4–6: Playwright executor + parser MVP (3 end-to-end flows)",
            "Week 7–8: Pass/fail reports + dashboard",
            "Week 9: Notify agent (team routing + ticket) — mentor feature",
            "Week 10–12: Batch suites, basic self-heal, polish, demo freeze",
            "Week 13–15: Report, poster, paper draft, final packaging",
        ],
    )

    # ---- SLIDE 8: Literature table ----
    s = prs.slides[7]
    for shape in s.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if full.strip().startswith("Literature"):
                set_shape_text(shape, ["Literature Review"])
        if shape.has_table:
            fill_table(
                shape.table,
                [
                    ["SOURCE / WORK", "YEAR", "APPROACH", "RESULT / INSIGHT", "RESEARCH GAPS"],
                    [
                        "Playwright Test Agents (Planner/Generator/Healer)",
                        "2025-26",
                        "Agentic plan → generate → heal tests",
                        "Industrial NL/plan-driven Playwright automation",
                        "Limited unified pass/fail docs + team notify workflow",
                    ],
                    [
                        "Tricentis Testim",
                        "Ongoing",
                        "AI smart locators / NL authoring",
                        "Reduces brittle UI automation maintenance",
                        "Closed commercial system; not open academic prototype",
                    ],
                    [
                        "Nass, Alegroth & Feldt (STVR)",
                        "2024",
                        "LLM for web element localization",
                        "Better localization when selectors fail",
                        "Not a full text-step executor + reporting pipeline",
                    ],
                    [
                        "ICECIT — Generative AI Self-Healing Selenium",
                        "2025",
                        "LLM vs rule-based locator repair",
                        "LLMs can outperform heuristics on repair",
                        "Focus on healing scripts, not NL execution + notify",
                    ],
                    [
                        "JAIGS — RL (PPO) Self-Heal in Playwright",
                        "2025",
                        "RL adaptive XPath recovery",
                        "Reduces maintenance under UI change",
                        "Does not cover scrum-style failure routing",
                    ],
                ],
            )

    # ---- SLIDE 9: Research gaps (5 pillars) ----
    s = prs.slides[8]
    gap_titles = {
        "Image-Centric": "Script Brittleness",
        "Accuracy Over": "Fragmented Toolchains",
        "Lack of Integrated": "Weak Evidence Trail",
        "Missing Stress": "No Ownership Routing",
        "Fragmented IDS": "Closed NL Platforms",
    }
    gap_bodies = {
        "Most adversarial robustness": "Traditional UI scripts break when selectors/DOM change, causing high maintenance cost in Quality Engineering.",
        "Existing IDS solutions": "Execution, reporting, and alerting are split across tools with no single agentic workflow.",
        "Very few intrusion": "Many runs lack structured step-level pass/fail documentation with screenshots as audit evidence.",
        "There is no standardized": "Failures are not automatically routed to the owning product team in a scrum-master style.",
        "Current approaches lack": "Commercial NL testing exists, but open explainable academic prototypes for industry PS are limited.",
    }
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Gap in the Research" in full or full.strip() in ("Gap Analysis", "Research Gaps"):
            set_shape_text(shape, ["Gap in the Research / Technology / Methodology"])
            continue
        replaced = False
        for k, v in gap_titles.items():
            if k in full and len(full) < 80:
                set_shape_text(shape, [v])
                replaced = True
                break
        if replaced:
            continue
        for k, v in gap_bodies.items():
            if k in full:
                set_shape_text(shape, [v])
                break

    # Also catch title if exact
    for shape in s.shapes:
        if shape.has_text_frame:
            t0 = shape.text_frame.paragraphs[0].text.strip()
            if "Gap" in t0 and len(t0) < 60:
                set_shape_text(shape, ["Gap in the Research / Technology / Methodology"])

    # ---- SLIDE 10: Proposed solution 3 pillars ----
    s = prs.slides[9]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if full.strip() == "Proposed Solution":
            set_shape_text(shape, ["Proposed Solution"])
        elif "Robust Machine Learning" in full:
            set_shape_text(shape, ["Text-to-Action Execution"])
        elif "Offline Adversarial" in full:
            set_shape_text(shape, ["Pass/Fail Documentation"])
        elif "Explainable and Deployable" in full:
            set_shape_text(shape, ["Scrum-Style Notify Agent"])
        elif "Develop an AI-driven intrusion" in full:
            set_shape_text(
                shape,
                [
                    "Parse plain-text steps into a locked JSON action schema and execute them on a web app using Playwright."
                ],
            )
        elif "Incorporate systematic offline" in full:
            set_shape_text(
                shape,
                [
                    "Produce step-level reports (status, expected vs actual, screenshots) so every run is auditable."
                ],
            )
        elif "Integrate SHAP-based" in full:
            set_shape_text(
                shape,
                [
                    "On failure, map module → owning team and raise an alert/ticket with failed-step context for fast triage."
                ],
            )

    # ---- SLIDE 11: Architecture - clear picture, keep title for user insert ----
    s = prs.slides[10]
    for shape in s.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if "Architectural" in full:
                set_shape_text(shape, ["Architectural Diagram"])
    clear_pictures(s)
    note = s.shapes.add_textbox(800000, 2800000, 7500000, 1200000)
    set_shape_text(
        note,
        [
            "[Insert architecture diagram here]",
            "Suggested flow: Text Steps → Parser → Playwright Executor → Assertions → Report Store → Notify Agent",
        ],
    )

    # ---- SLIDE 12: Technology stack ----
    s = prs.slides[11]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if full.strip() == "Technology Stack":
            set_shape_text(shape, ["Technology Stack"])
        elif "Machine Learning" in full or "Scikit-learn" in full or "Development Environment" in full:
            set_shape_text(
                shape,
                [
                    "1. Browser Automation",
                    "Playwright (Chromium) — action execution, assertions, screenshots",
                    "",
                    "2. Agent / AI Layer",
                    "Rule-based parser (Phase-1) + LLM parser (Phase-2) for text → JSON steps",
                    "Pydantic schemas for validation",
                    "",
                    "3. Backend & Reporting",
                    "Python, FastAPI — run APIs",
                    "JSON + Markdown report writers",
                    "YAML team ownership map",
                    "",
                    "4. Notification",
                    "Console alerts (now); Slack webhook / Email (next)",
                    "",
                    "5. Engineering",
                    "GitHub, pytest, VS Code / Cursor",
                    "Demo apps: Sauce Demo / DemoQA (ClickUp/Zoho if mentors allow)",
                ],
            )

    # ---- SLIDE 13: Work done 1 ----
    s = prs.slides[12]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Work done" in full and len(full) < 40:
            set_shape_text(shape, ["Work done till now"])
        elif "Data Preparation" in full or "UNSW" in full or "Preprocessing" in full:
            set_shape_text(
                shape,
                [
                    "Foundation & Contracts:",
                    "Created project repository and package layout",
                    "Locked Step Schema (goto/fill/click/assert/...)",
                    "Locked Report Schema (pass/fail + evidence fields)",
                    "Team ownership map (module → team → channel)",
                    "Wrote objectives, architecture notes, literature starter",
                    "Prepared 10 plain-text sample test cases",
                    "Output",
                    "Runnable scaffold with fixed contracts for integration",
                ],
            )

    # ---- SLIDE 14: Work done 2 ----
    s = prs.slides[13]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Work done" in full and len(full) < 40:
            set_shape_text(shape, ["Work done till now"])
        elif "Model Development" in full or "Baseline Models" in full or "Logistic Regression" in full:
            set_shape_text(
                shape,
                [
                    "Core Pipeline Implementation:",
                    "Rule-based plain-text → JSON parser",
                    "PlaywrightExecutor (actions + assertions)",
                    "Failure screenshots + skip-after-fail policy",
                    "Report writer (JSON + Markdown)",
                    "NotifyAgent (console ticket-style alerts)",
                    "FastAPI skeleton (/run/text, /run/json)",
                    "CLI runner scripts/run_suite.py",
                    "Outcome",
                    "End-to-end executable MVP path available",
                ],
            )

    # ---- SLIDE 15: Work done 3 ----
    s = prs.slides[14]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Work done" in full and len(full) < 40:
            set_shape_text(shape, ["Work done till now"])
        elif "Evaluation" in full or "Performance Evaluation" in full or "SHAP" in full:
            set_shape_text(
                shape,
                [
                    "Validation & Smoke Results:",
                    "Unit test for parser mapping (pytest)",
                    "JSON login suite on Sauce Demo: PASSED (7/7 steps)",
                    "Intentional fail case: FAILED at assert_text as expected",
                    "Notify agent triggered ticket to mapped QA team",
                    "Artifacts stored under data/reports and screenshots",
                    "Next",
                    "LLM parser, dashboard UI, more flows, Slack/email notify",
                ],
            )

    # ---- SLIDE 16: Conclusion ----
    s = prs.slides[15]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if "Conclusion" in full and len(full) < 50:
            set_shape_text(shape, ["Conclusion"])
        elif "Tree-Based" in full or "XGBoost" in full or "overfitting" in full:
            set_shape_text(
                shape,
                [
                    "DS 1 needs an agentic Quality Engineering workflow: text steps → execute → verify.",
                    "Pass/fail documentation with evidence is essential for mentor/industry review.",
                    "Scrum-style notify closes the loop by routing failures to owning teams quickly.",
                    "Phase-1 foundation is demoable; Phase-2 will deepen AI parsing, UI, and evaluation.",
                ],
            )

    # ---- SLIDE 17: References ----
    s = prs.slides[16]
    # Collect text boxes in order roughly by top
    ref_shapes = []
    for shape in s.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if full.strip() == "References":
                set_shape_text(shape, ["References"])
            elif "http" in full or "Explaining" in full or "Towards" in full or "Practical" in full or "Robustness" in full:
                ref_shapes.append(shape)
    ref_shapes.sort(key=lambda sh: sh.top)
    new_refs = [
        ["https://playwright.dev/docs/test-agents", "Playwright Test Agents (Planner, Generator, Healer) — Microsoft"],
        ["https://www.testim.io/", "Tricentis Testim — AI-powered web test automation (named in Dassault brief)"],
        ["https://doi.org/10.1002/stvr.1893", "Nass et al. (2024) — Improving Web Element Localization Using an LLM (STVR)"],
        ["https://doi.org/10.1109/icecit67774.2025.11451150", "ICECIT 2025 — Generative AI for Self-Healing Selenium Tests"],
        ["https://doi.org/10.60087/jaigs.v5i1.341", "JAIGS 2025 — Self-Healing Automation with RL (PPO) in Playwright"],
    ]
    for shape, (url, title) in zip(ref_shapes, new_refs):
        set_shape_text(shape, [url, title])

    # ---- SLIDE 18: photograph ----
    s = prs.slides[17]
    for shape in s.shapes:
        if shape.has_text_frame:
            full = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if "photograph" in full.lower():
                set_shape_text(shape, ["A photograph of the presentation after due permission from the guide."])

    replace_all_dates(prs)
    prs.save(OUT)
    shutil.copy2(OUT, OUT_DL)
    print("Saved:", OUT)
    print("Copied:", OUT_DL)
    print("Slides:", len(prs.slides))


if __name__ == "__main__":
    main()
