"""Create DS1 Gantt timeline image and insert into CA1 PPT slide 7."""

from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

PPT = Path(r"C:\Users\shwet\Downloads\DS1_CA1_Presentation_from_PBL_Template.pptx")
OUT_PPT = PPT
OUT_PROJ = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\DS1_CA1_Presentation_from_PBL_Template.pptx")
IMG = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\ca1_assets\ds1_timeline_gantt.png")

# Colors matching the reference image
HEADER_BG = (47, 64, 80)       # dark slate
ROW_LABEL_BG = (232, 232, 232) # light gray
WHITE = (255, 255, 255)
BAR = (36, 52, 71)             # navy bars
GRID = (255, 255, 255)
TEXT = (20, 20, 20)
HEADER_TEXT = (255, 255, 255)

# DS1 phases mapped onto 12 timeline units (same visual format as template)
ROWS = [
    (
        "Ideation and Research",
        "Finalize DS1 concept, literature review,\nstudy Playwright/agentic testing approaches.",
        1,
        2,
    ),
    (
        "Design & Contracts",
        "Step/report schemas, architecture freeze,\nsample text cases, ownership map.",
        2,
        4,
    ),
    (
        "Core Executor MVP",
        "Playwright actions/assertions, text→JSON\nparser, end-to-end flow integration.",
        5,
        7,
    ),
    (
        "Reporting & Dashboard",
        "Pass/fail documentation with evidence,\nrun history dashboard and exports.",
        7,
        8,
    ),
    (
        "Notify & Enhancements",
        "Scrum-style failure notify, batch suites,\nbasic self-heal / polish.",
        8,
        10,
    ),
    (
        "Documentation and Review",
        "Report/paper preparation, demo freeze,\nposter and final review packaging.",
        11,
        12,
    ),
]


def font(size: int, bold: bool = False):
    candidates = [
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap_fit(draw, text, font_obj, max_width):
    # Keep explicit newlines; wrap long lines
    out_lines = []
    for raw in text.split("\n"):
        words = raw.split()
        cur = ""
        for w in words:
            trial = (cur + " " + w).strip()
            if draw.textlength(trial, font=font_obj) <= max_width:
                cur = trial
            else:
                if cur:
                    out_lines.append(cur)
                cur = w
        if cur:
            out_lines.append(cur)
    return out_lines


def make_gantt(path: Path):
    # High-res for PPT clarity
    W, H = 2200, 980
    img = Image.new("RGB", (W, H), WHITE)
    draw = ImageDraw.Draw(img)

    margin = 20
    header_h = 70
    # Column widths
    init_w = 420
    obj_w = 720
    timeline_w = W - margin * 2 - init_w - obj_w
    col_w = timeline_w / 12
    row_h = (H - margin * 2 - header_h) / len(ROWS)

    x0 = margin
    y0 = margin

    f_head = font(28, bold=True)
    f_cell = font(22, bold=True)
    f_obj = font(20, bold=False)
    f_num = font(22, bold=True)

    # Header background
    draw.rectangle([x0, y0, W - margin, y0 + header_h], fill=HEADER_BG)

    # Header labels
    def center_text(text, box, fnt, fill):
        x1, y1, x2, y2 = box
        bbox = draw.textbbox((0, 0), text, font=fnt)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        draw.text((x1 + (x2 - x1 - tw) / 2, y1 + (y2 - y1 - th) / 2), text, font=fnt, fill=fill)

    center_text("Initiative", (x0, y0, x0 + init_w, y0 + header_h), f_head, HEADER_TEXT)
    center_text("Objective", (x0 + init_w, y0, x0 + init_w + obj_w, y0 + header_h), f_head, HEADER_TEXT)
    for i in range(12):
        cx1 = x0 + init_w + obj_w + i * col_w
        cx2 = cx1 + col_w
        center_text(f"{i+1:02d}", (cx1, y0, cx2, y0 + header_h), f_num, HEADER_TEXT)

    # Header vertical separators (subtle)
    for x in [x0 + init_w, x0 + init_w + obj_w]:
        draw.line([x, y0, x, y0 + header_h], fill=(70, 90, 110), width=2)
    for i in range(1, 12):
        x = x0 + init_w + obj_w + i * col_w
        draw.line([x, y0, x, y0 + header_h], fill=(70, 90, 110), width=1)

    # Rows
    for r, (initiative, objective, start, end) in enumerate(ROWS):
        y1 = y0 + header_h + r * row_h
        y2 = y1 + row_h

        # label cols
        draw.rectangle([x0, y1, x0 + init_w, y2], fill=ROW_LABEL_BG)
        draw.rectangle([x0 + init_w, y1, x0 + init_w + obj_w, y2], fill=ROW_LABEL_BG)
        # timeline area white
        draw.rectangle([x0 + init_w + obj_w, y1, W - margin, y2], fill=WHITE)

        # initiative text
        lines = wrap_fit(draw, initiative, f_cell, init_w - 30)
        total_h = len(lines) * 28
        ty = y1 + (row_h - total_h) / 2
        for li, line in enumerate(lines):
            bbox = draw.textbbox((0, 0), line, font=f_cell)
            tw = bbox[2] - bbox[0]
            draw.text((x0 + (init_w - tw) / 2, ty + li * 28), line, font=f_cell, fill=TEXT)

        # objective text (left aligned with padding)
        olines = wrap_fit(draw, objective, f_obj, obj_w - 36)
        total_h = len(olines) * 24
        ty = y1 + (row_h - total_h) / 2
        for li, line in enumerate(olines):
            draw.text((x0 + init_w + 16, ty + li * 24), line, font=f_obj, fill=TEXT)

        # grid lines for timeline cells
        for i in range(13):
            x = x0 + init_w + obj_w + i * col_w
            draw.line([x, y1, x, y2], fill=(230, 230, 230), width=1)
        draw.line([x0 + init_w + obj_w, y2, W - margin, y2], fill=(230, 230, 230), width=1)

        # bar from start to end inclusive
        bx1 = x0 + init_w + obj_w + (start - 1) * col_w + 6
        bx2 = x0 + init_w + obj_w + end * col_w - 6
        by1 = y1 + row_h * 0.28
        by2 = y2 - row_h * 0.28
        draw.rounded_rectangle([bx1, by1, bx2, by2], radius=8, fill=BAR)

        # outer borders for label cells
        draw.rectangle([x0, y1, x0 + init_w, y2], outline=(200, 200, 200), width=1)
        draw.rectangle([x0 + init_w, y1, x0 + init_w + obj_w, y2], outline=(200, 200, 200), width=1)

    # Outer border
    draw.rectangle([x0, y0, W - margin, H - margin], outline=(180, 180, 180), width=2)
    # vertical separators full height
    draw.line([x0 + init_w, y0, x0 + init_w, H - margin], fill=(180, 180, 180), width=2)
    draw.line([x0 + init_w + obj_w, y0, x0 + init_w + obj_w, H - margin], fill=(180, 180, 180), width=2)

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG")
    print("Saved image:", path)


def clear_non_title_content(slide):
    """Remove old timeline text boxes/pictures but keep title + footer-ish shapes if possible."""
    keep_keywords = ("project plan", "timeline", "department of artificial", "28-07-2026", "17-01-2025")
    to_remove = []
    for shape in slide.shapes:
        # Always remove pictures on this slide (old gantt / placeholders)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            to_remove.append(shape)
            continue
        if not shape.has_text_frame:
            # decorative groups/side bars - keep
            continue
        text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip().lower()
        if not text:
            continue
        # keep title
        if "project plan" in text or text in {"7", "8"} or "department of artificial" in text or "28-07-2026" in text:
            continue
        # remove previous bullet timeline text we added
        if "week 1" in text or "kickoff" in text or "insert architecture" in text or text.startswith("week "):
            to_remove.append(shape)
            continue
        # remove any other large body text boxes
        if shape.width > Inches_safe(5) and shape.top > Inches_safe(1.2):
            to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)


def Inches_safe(n):
    return Inches(n)


def insert_into_ppt():
    make_gantt(IMG)
    prs = Presentation(str(PPT if PPT.exists() else OUT_PROJ))
    slide = prs.slides[6]  # Project plan with timeline

    # Ensure title is correct
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if "Project plan" in t or "timeline" in t.lower():
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = "Project plan with timeline"
                break

    clear_non_title_content(slide)

    # Insert image - similar placement to original template picture
    # Original was approx L=177940 T=1555565 W=8776355 H=4180919
    left = Emu(177940)
    top = Emu(1555565)
    width = Emu(8776355)
    slide.shapes.add_picture(str(IMG), left, top, width=width)

    prs.save(str(OUT_PROJ))
    shutil.copy2(OUT_PROJ, OUT_PPT)
    print("Updated PPT:", OUT_PPT)


if __name__ == "__main__":
    insert_into_ppt()
