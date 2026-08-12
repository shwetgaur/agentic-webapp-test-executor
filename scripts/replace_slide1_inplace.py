"""Rebuild slide 1 in-place to match previous B.Tech CA1 title slide."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PREV = Path(r"C:\Users\shwet\Downloads\DS1_Agentic_WebApp_Test_Executor_CA1_Presentation.pptx")
# Use project copy of filled PBL deck if Downloads was corrupted; regenerate from template fill output
SRC_GOOD = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\DS1_CA1_Presentation_from_PBL_Template.pptx")
TEMPLATE3 = Path(r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (3).pptx")
OUT = Path(r"C:\Users\shwet\Downloads\DS1_CA1_Presentation_from_PBL_Template.pptx")
OUT_PROJ = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\DS1_CA1_Presentation_from_PBL_Template.pptx")
FILL_SCRIPT_DONE_MARKER = OUT_PROJ


def clear_slide_shapes(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_text(slide, left, top, width, height, text, size=18, bold=False, align=PP_ALIGN.LEFT, color=RGBColor(0x1A, 0x1A, 0x1A)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def extract_prev_texts_and_logo():
    prev = Presentation(str(PREV))
    texts = {}
    logo_blob = None
    for sh in prev.slides[0].shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            logo_blob = sh.image.blob
            continue
        if not sh.has_text_frame:
            continue
        t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
        if t.startswith("B. Tech"):
            texts["header"] = t
        elif t.startswith("CA 1"):
            texts["ca"] = t
        elif t.startswith("Project ID"):
            texts["pid"] = t
        elif t.startswith("Title"):
            texts["title"] = t
        elif t.startswith("Name of the Guide"):
            texts["guide"] = t
        elif t.startswith("Name of the Industry"):
            texts["mentor"] = t
        elif t.startswith("Group Members"):
            texts["members"] = t
        elif "Department" in t:
            texts["dept"] = t
    return texts, logo_blob


def main():
    # Re-run content fill from original template to get a clean 18-slide file, then only change slide 1.
    # Import fill logic by re-executing fill script output path then patching slide 1.
    import subprocess
    import sys

    subprocess.check_call([sys.executable, str(Path(__file__).with_name("fill_pbl_template_ds1.py"))])

    # fill script overwrote with old member names - that's ok, we replace entire slide 1 from previous
    prs = Presentation(str(OUT_PROJ))
    slide = prs.slides[0]
    clear_slide_shapes(slide)

    texts, logo_blob = extract_prev_texts_and_logo()

    add_text(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4), texts.get("header", "B. Tech. Project AY 2026-27"), size=28, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), texts.get("ca", "CA 1 Presentation"), size=22, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(1.15), Inches(9), Inches(0.3), texts.get("pid", "Project ID: DS 1"), size=16, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), Inches(1.5), Inches(9.2), Inches(0.55), texts.get("title", "Title of the Project: Agentic Web-App Test Executor\n(Domain: Quality Engineering)"), size=16, bold=True, align=PP_ALIGN.CENTER)

    tmp_logo = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\ca1_assets\_tmp_title_logo.png")
    if logo_blob:
        tmp_logo.write_bytes(logo_blob)
        slide.shapes.add_picture(str(tmp_logo), Inches(2.7), Inches(2.15), width=Inches(4.6))
    else:
        fallback = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\ca1_assets\sit_logo_title.png")
        if fallback.exists():
            slide.shapes.add_picture(str(fallback), Inches(2.7), Inches(2.15), width=Inches(4.6))

    add_text(slide, Inches(0.5), Inches(4.35), Inches(9), Inches(0.3), texts.get("guide", "Name of the Guide: Mayur Gaikwad"), size=15, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(0.3), texts.get("mentor", "Name of the Industry Mentor: Dassault Systemes (ENOVIA) / ________________"), size=14, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.15), Inches(8.5), Inches(1.55), texts.get("members", "Group Members:"), size=14)

    gray = RGBColor(0x66, 0x66, 0x66)
    add_text(slide, Inches(0.35), Inches(7.05), Inches(2.2), Inches(0.3), "28-07-2026", size=11, color=gray)
    add_text(slide, Inches(2.4), Inches(7.05), Inches(5.5), Inches(0.3), texts.get("dept", "Department of Artificial Intelligence & Machine Learning"), size=11, align=PP_ALIGN.CENTER, color=gray)
    add_text(slide, Inches(8.7), Inches(7.05), Inches(0.9), Inches(0.3), "1", size=11, align=PP_ALIGN.RIGHT, color=gray)

    # Also update work-distribution / other slides that still have old names if fill script reset them
    # Sync member names into slide 3 table from previous title members
    members_block = texts.get("members", "")
    # Parse names from members block for table if possible
    # Keep work distribution roles mapped by order
    name_lines = [ln.strip() for ln in members_block.splitlines() if ln.strip() and not ln.strip().startswith("Group")]
    # lines like "1. Shwet Gaur (PRN: ...)"
    names = []
    for ln in name_lines:
        # remove leading number.
        part = ln.split(".", 1)[-1].strip()
        name = part.split("(PRN")[0].split("(prn")[0].strip()
        names.append(name)

    roles = [
        "End-to-end pipeline ownership: step schema, parser, Playwright executor, notify agent, integration.",
        "Playwright action coverage, assertions, demo scenarios, failure screenshot tooling.",
        "Platform APIs/DB, dashboard & report storage, notify service plumbing, GitHub structure.",
        "QA sample test cases, pass/fail report review, literature support, evaluation notes.",
    ]
    # Map: user listed Shwet, Sahishnu, Eesha, Saksham - assign sensible roles
    if len(names) >= 4:
        for shape in prs.slides[2].shapes:
            if shape.has_table:
                table = shape.table
                # header remains
                for i in range(4):
                    table.cell(i + 1, 0).text = names[i]
                    table.cell(i + 1, 1).text = roles[i]

    prs.save(str(OUT_PROJ))
    shutil.copy2(OUT_PROJ, OUT)
    print("OK slide1 replaced in-place")
    print(OUT)
    print("slides", len(prs.slides))
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            t = " | ".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
            if t:
                print(t[:130])


if __name__ == "__main__":
    main()
