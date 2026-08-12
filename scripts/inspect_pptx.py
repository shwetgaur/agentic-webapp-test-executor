from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

paths = [
    r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (2).pptx",
    r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (1).pptx",
    r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA.pptx",
]

for path in paths:
    print("\n" + "=" * 80)
    print(path)
    prs = Presentation(path)
    print("size", prs.slide_width, prs.slide_height, "slides", len(prs.slides))
    for i, layout in enumerate(prs.slide_layouts):
        print(f"LAYOUT {i}: {layout.name}")
    for si, slide in enumerate(prs.slides[:8]):
        print(f"\n----- SLIDE {si+1} -----")
        for shape in slide.shapes:
            texts = []
            if shape.has_text_frame:
                texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            print(
                {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "l": shape.left,
                    "t": shape.top,
                    "w": shape.width,
                    "h": shape.height,
                    "text": texts[:8],
                }
            )
