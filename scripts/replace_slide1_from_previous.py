"""Replace slide 1 of PBL-based DS1 PPT with title slide from previous B.Tech CA1 PPT."""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

PREV = Path(r"C:\Users\shwet\Downloads\DS1_Agentic_WebApp_Test_Executor_CA1_Presentation.pptx")
CUR = Path(r"C:\Users\shwet\Downloads\DS1_CA1_Presentation_from_PBL_Template.pptx")
OUT = Path(r"C:\Users\shwet\Downloads\DS1_CA1_Presentation_from_PBL_Template.pptx")
OUT_PROJ = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\DS1_CA1_Presentation_from_PBL_Template.pptx")


def delete_slide(prs: Presentation, index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def append_slide_from(src_prs: Presentation, src_index: int, dst_prs: Presentation) -> None:
    """Append a deep-copied slide from src into dst (end), preserving shapes via XML clone of slide part is complex;
    instead: add blank slide and copy all shape XML from source slide.
    """
    src_slide = src_prs.slides[src_index]
    # Use blank layout
    blank = dst_prs.slide_layouts[6] if len(dst_prs.slide_layouts) > 6 else dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(blank)

    # Remove any default shapes on blank
    for shape in list(dst_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy each shape element from source
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        dst_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    # Copy slide-level image parts: pictures need related image blobs
    # Re-add pictures properly by checking picture shapes
    # If pictures fail to resolve, recreate title slide from scratch using assets.


def move_last_slide_to_front(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    last = items[-1]
    sldIdLst.remove(last)
    sldIdLst.insert(0, last)


def rebuild_title_slide_like_previous(dst_prs: Presentation) -> None:
    """Fallback: rebuild slide 1 visually like previous CA1 title slide."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt, Emu

    # Delete current first slide and create new blank at end then move
    delete_slide(dst_prs, 0)
    blank = dst_prs.slide_layouts[6]
    s = dst_prs.slides.add_slide(blank)
    move_last_slide_to_front(dst_prs)
    s = dst_prs.slides[0]

    # Clear leftovers
    for shape in list(s.shapes):
        shape._element.getparent().remove(shape._element)

    def add_text(left, top, width, height, text, size=18, bold=False, align=PP_ALIGN.LEFT, color=RGBColor(0x1A, 0x1A, 0x1A)):
        box = s.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        # multi-line support
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"
        return box

    # Match previous slide content (already updated by user)
    prev = Presentation(PREV)
    texts = {}
    logo_blob = None
    for sh in prev.slides[0].shapes:
        if sh.has_text_frame:
            t = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
            if t.startswith("B. Tech"):
                texts["header"] = t
            elif t.startswith("CA 1"):
                texts["ca"] = t
            elif t.startswith("Project ID"):
                texts["pid"] = t
            elif t.startswith("Title"):
                texts["title"] = t
            elif t.startswith("Name of the Guide"):
                texts["guide"] = t
            elif t.startswith("Name of the Industry"):
                texts["mentor"] = t
            elif t.startswith("Group Members"):
                texts["members"] = t
            elif "Department" in t:
                texts["dept"] = t
            elif t.strip() in ("1", "28-07-2026", "20-07-2026"):
                pass
        if sh.shape_type is not None and str(sh.shape_type) == "PICTURE (13)":
            logo_blob = sh.image.blob

    add_text(Inches(0.5), Inches(0.25), Inches(9), Inches(0.4), texts.get("header", "B. Tech. Project AY 2026-27"), size=28, bold=True, align=PP_ALIGN.CENTER)
    add_text(Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), texts.get("ca", "CA 1 Presentation"), size=22, bold=True, align=PP_ALIGN.CENTER)
    add_text(Inches(0.5), Inches(1.15), Inches(9), Inches(0.3), texts.get("pid", "Project ID: DS 1"), size=16, align=PP_ALIGN.CENTER)
    add_text(Inches(0.4), Inches(1.5), Inches(9.2), Inches(0.55), texts.get("title", "Title of the Project: Agentic Web-App Test Executor\n(Domain: Quality Engineering)"), size=16, bold=True, align=PP_ALIGN.CENTER)

    # logo
    logo_path = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\ca1_assets\sit_logo_title.png")
    if logo_blob:
        tmp = Path(r"C:\Users\shwet\Projects\agentic-webapp-test-executor\docs\ca1_assets\_tmp_title_logo.png")
        tmp.write_bytes(logo_blob)
        s.shapes.add_picture(str(tmp), Inches(2.7), Inches(2.15), width=Inches(4.6))
    elif logo_path.exists():
        s.shapes.add_picture(str(logo_path), Inches(2.7), Inches(2.15), width=Inches(4.6))

    add_text(Inches(0.5), Inches(4.35), Inches(9), Inches(0.3), texts.get("guide", "Name of the Guide: Mayur Gaikwad"), size=15, align=PP_ALIGN.CENTER)
    add_text(Inches(0.5), Inches(4.7), Inches(9), Inches(0.3), texts.get("mentor", "Name of the Industry Mentor: Dassault Systemes (ENOVIA) / ________________"), size=14, align=PP_ALIGN.CENTER)
    add_text(Inches(0.6), Inches(5.15), Inches(8.5), Inches(1.5), texts.get("members", "Group Members:"), size=14)

    # footer
    add_text(Inches(0.35), Inches(7.05), Inches(2.2), Inches(0.3), "28-07-2026", size=11, color=RGBColor(0x66, 0x66, 0x66))
    add_text(Inches(2.4), Inches(7.05), Inches(5.5), Inches(0.3), texts.get("dept", "Department of Artificial Intelligence & Machine Learning"), size=11, align=PP_ALIGN.CENTER, color=RGBColor(0x66, 0x66, 0x66))
    add_text(Inches(8.7), Inches(7.05), Inches(0.9), Inches(0.3), "1", size=11, align=PP_ALIGN.RIGHT, color=RGBColor(0x66, 0x66, 0x66))


def main():
    # Work on a copy in memory from current PBL deck
    shutil.copy2(CUR, OUT_PROJ)
    dst = Presentation(str(OUT_PROJ))
    rebuild_title_slide_like_previous(dst)
    dst.save(str(OUT_PROJ))
    shutil.copy2(OUT_PROJ, OUT)
    print("Updated slide 1 from previous CA1 title slide")
    print("Saved:", OUT)

    # verify
    prs = Presentation(str(OUT))
    print("Slide count:", len(prs.slides))
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            t = " | ".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
            if t:
                print(" ", t[:140])


if __name__ == "__main__":
    main()
