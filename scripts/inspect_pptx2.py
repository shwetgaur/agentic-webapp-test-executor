from pptx import Presentation

path = r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (2).pptx"
prs = Presentation(path)
print("slides", len(prs.slides))
for si, slide in enumerate(prs.slides):
    print(f"\n===== SLIDE {si+1} =====")
    for shape in slide.shapes:
        texts = []
        if shape.has_text_frame:
            texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
        has_pic = shape.shape_type is not None and str(shape.shape_type) == "PICTURE (13)"
        print(shape.name, shape.shape_type, texts[:6], "PIC" if has_pic else "")
