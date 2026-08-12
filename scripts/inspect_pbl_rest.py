from pptx import Presentation
import sys
sys.stdout.reconfigure(encoding="utf-8")

path = r"C:\Users\shwet\Downloads\PBL-2 Review Presentation (CA-1)Template_CA (3).pptx"
prs = Presentation(path)
for si in range(15, len(prs.slides)):
    slide = prs.slides[si]
    print(f"\n===== SLIDE {si+1} =====")
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                print(shape.name, "->")
                for t in texts:
                    print(" ", t[:200])
        if shape.has_table:
            print("TABLE")
            for r in shape.table.rows:
                print(" | ".join(c.text.replace("\n"," / ")[:80] for c in r.cells))
